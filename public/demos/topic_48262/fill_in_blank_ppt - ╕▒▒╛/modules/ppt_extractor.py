"""PPT挖空文本提取模块

从已生成的挖空PPT中逆向提取出带 {{}} 标记的文本文件。

原理：
1. 基础层（placeholder）：普通文字 + noFill/uFill（透明文字+下划线）的挖空词
2. 答案层（无placeholder的shape）：蓝色文字（accent1+lumMod+lumOff）为答案
3. 通过对比基础层的下划线文字和答案层的蓝色文字，还原 {{}} 标记
"""
import os
import re
from typing import List, Tuple
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree


def extract_blanks_from_pptx(pptx_path: str, output_txt_path: str = None) -> str:
    """从挖空PPT中提取带 {{}} 标记的文本

    Args:
        pptx_path: PPT文件路径
        output_txt_path: 输出文本文件路径（可选，不指定则自动生成）

    Returns:
        提取的文本内容（带 {{}} 标记）
    """
    prs = Presentation(pptx_path)

    # 按 slide 顺序提取，根据标题页/内容页重组主题
    topics = []  # [(topic_name, [paragraphs]), ...]
    current_topic = None
    current_paragraphs = []

    for slide_idx, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)

        # 找标题和内容
        title_text = ""
        base_paras = []  # [(para_idx, [(text, is_blank), ...]), ...]
        answer_blues = {}  # {shape_id: [(para_idx, blue_text), ...]}

        for shape in shapes:
            if not shape.has_text_frame:
                continue

            sp = shape._element
            cNvPr = sp.find('.//' + qn('p:cNvPr'))
            sid = cNvPr.get('id') if cNvPr is not None else '?'
            ph = sp.find('.//' + qn('p:ph'))

            txBody = sp.find('.//' + qn('p:txBody'))
            if txBody is None:
                continue

            paras = txBody.findall(qn('a:p'))

            # 标题（type="title" 或 type="ctrTitle"）
            if ph is not None:
                ph_type = ph.get('type') or ''
                if ph_type in ('title', 'ctrTitle'):
                    for p in paras:
                        runs = p.findall(qn('a:r'))
                        texts = [r.find(qn('a:t')).text or '' for r in runs if r.find(qn('a:t')) is not None]
                        title_text = ''.join(texts).strip()
                    continue

                # 基础层（ph idx="1"）
                ph_idx = ph.get('idx')
                if ph_idx == '1':
                    for pi, p in enumerate(paras):
                        runs = p.findall(qn('a:r'))
                        segments = []
                        for r in runs:
                            t_el = r.find(qn('a:t'))
                            text = t_el.text if t_el is not None else ''
                            rPr = r.find(qn('a:rPr'))
                            is_blank = False
                            if rPr is not None:
                                # noFill + uFill = 挖空词
                                noFill = rPr.find(qn('a:noFill'))
                                uFill = rPr.find(qn('a:uFill'))
                                u_attr = rPr.get('u')
                                if (noFill is not None and uFill is not None) or u_attr == 'sng':
                                    is_blank = True
                            if text:
                                segments.append((text, is_blank))
                        if segments:
                            base_paras.append((pi, segments))
                    continue

            # 答案层（无 ph）
            # 收集蓝色文字
            for pi, p in enumerate(paras):
                runs = p.findall(qn('a:r'))
                for r in runs:
                    t_el = r.find(qn('a:t'))
                    text = t_el.text if t_el is not None else ''
                    rPr = r.find(qn('a:rPr'))
                    if rPr is not None:
                        sf = rPr.find(qn('a:solidFill'))
                        if sf is not None:
                            sc = sf.find(qn('a:schemeClr'))
                            if sc is not None and sc.get('val') == 'accent1':
                                # 蓝色答案文字
                                if sid not in answer_blues:
                                    answer_blues[sid] = []
                                answer_blues[sid].append((pi, text))

        # 如果有标题，检查是否是新主题（避免分页导致的重复标题）
        if title_text:
            if title_text != current_topic:
                # 新主题
                if current_topic and current_paragraphs:
                    topics.append((current_topic, current_paragraphs))
                current_topic = title_text
                current_paragraphs = []
            # 如果标题与当前主题相同（分页），不重复添加标题，只追加段落

        # 处理基础层段落
        if base_paras:
            for pi, segments in base_paras:
                para_text = ""
                skip_next = False
                prev_blank_text = ""
                for text, is_blank in segments:
                    if skip_next and text == prev_blank_text:
                        # 跳过与上一个 blank 同名的 non-blank run
                        skip_next = False
                        continue
                    skip_next = False
                    if is_blank:
                        para_text += "{{" + text + "}}"
                        prev_blank_text = text
                        skip_next = True  # 标记跳过下一个同名 run
                    else:
                        para_text += text
                if para_text.strip():
                    current_paragraphs.append(para_text)

    # 最后一个主题
    if current_topic and current_paragraphs:
        topics.append((current_topic, current_paragraphs))

    # 生成文本
    lines = []
    for topic_name, paragraphs in topics:
        lines.append(f"%{topic_name}")
        for para in paragraphs:
            lines.append(para)
            lines.append("")  # 空行分隔段落
        lines.append("")  # 额外空行分隔主题

    content = '\n'.join(lines).strip() + '\n'

    # 写入文件
    if output_txt_path is None:
        output_txt_path = os.path.splitext(pptx_path)[0] + '_extracted.txt'

    with open(output_txt_path, 'w', encoding='utf-8') as f:
        f.write(content)

    print(f"[OK] Extracted text saved to: {output_txt_path}")
    print(f"     Topics: {len(topics)}")
    total_blanks = content.count('{{')
    print(f"     Blanks: {total_blanks}")

    return content


if __name__ == '__main__':
    import sys
    if len(sys.argv) < 2:
        print("Usage: python ppt_extractor.py <input.pptx> [output.txt]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    extract_blanks_from_pptx(pptx_path, output_path)

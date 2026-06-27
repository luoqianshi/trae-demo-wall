"""PPT生成模块 v2 - 基于模板占位符克隆，完全匹配参考PPT结构

核心改变：
1. 不再用 add_textbox/add_shape 从零构建
2. 直接从模板的 slideLayout 复制占位符 shape
3. 克隆时保留 spLocks、lstStyle、p:ph 等所有继承属性
4. 坐标使用整数 EMU（与 layout 完全对齐）
"""
import os
import copy
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from .input_parser import Topic, Paragraph
from .xml_utils import (
    add_click_fade_animation,
    split_text_by_blanks,
)
from .text_paginator import paginate_topics, TextMeasurer

# 整数 EMU 常量（与参考PPT完全一致）
SLIDE_W = 12192000  # 13.33 inch
SLIDE_H = 6858000   # 7.50 inch

# 内容区坐标（整数 EMU，从参考PPT提取）
CONTENT_LEFT = 608400    # 0.6654 inch (近似)
CONTENT_TOP = 1490400    # 1.6299 inch (近似)
CONTENT_W = 10969200     # 11.9961 inch
CONTENT_H = 4759200      # 5.2047 inch

# 标题区坐标
TITLE_LEFT = 608400
TITLE_TOP = 608400       # 0.6654 inch
TITLE_W = 10969200
TITLE_H = 707600         # 0.7717 inch

# 标题页图片坐标
IMG_LEFT = 9820200       # 10.75 inch
IMG_TOP = 5168900        # 5.66 inch
IMG_W = 2356000          # 2.58 inch
IMG_H = 1680300          # 1.84 inch

# 边距（整数 EMU）
MARGIN_L = 90000
MARGIN_R = 90000
MARGIN_T = 46800
MARGIN_B = 46800

FONT_SIZE = 3200  # 32pt = 3200 (in 1/100 pt)
LINE_SPACING = 130000  # 130% = 130000 (in 1/1000 %)
SPACE_AFTER = 100  # 1pt = 100 (in 1/100 pt)

# 答案层shape的完整lstStyle（从参考PPT提取，lvl1pPr-lvl9pPr）
# 答案shape没有p:ph，无法从layout继承，必须自带lstStyle
# 关键：lvl1pPr包含 lnSpc=130000(130%), spcBef=0, spcAft=1000(10pt)
ANSWER_LST_STYLE = """<a:lstStyle>
    <a:lvl1pPr marL="228600" indent="-228600" algn="l" defTabSz="914400" rtl="0" eaLnBrk="1" fontAlgn="auto" latinLnBrk="0" hangingPunct="1">
      <a:lnSpc><a:spcPct val="130000"/></a:lnSpc>
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="1000"/></a:spcAft>
      <a:buFont typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>
      <a:buChar char="●"/>
      <a:defRPr sz="1800" u="none" strike="noStrike" kern="1200" cap="none" spc="150" normalizeH="0" baseline="0">
        <a:solidFill><a:schemeClr val="tx1"><a:lumMod val="65000"/><a:lumOff val="35000"/></a:schemeClr></a:solidFill>
        <a:uFillTx/>
        <a:latin typeface="+mn-lt"/>
        <a:ea typeface="+mn-ea"/>
        <a:cs typeface="+mn-cs"/>
      </a:defRPr>
    </a:lvl1pPr>
    <a:lvl2pPr marL="685800" indent="-228600" algn="l" defTabSz="914400" rtl="0" eaLnBrk="1" fontAlgn="auto" latinLnBrk="0" hangingPunct="1">
      <a:lnSpc><a:spcPct val="120000"/></a:lnSpc>
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="600"/></a:spcAft>
      <a:buFont typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>
      <a:buChar char="●"/>
      <a:tabLst><a:tab pos="1609725" algn="l"/><a:tab pos="1609725" algn="l"/><a:tab pos="1609725" algn="l"/><a:tab pos="1609725" algn="l"/></a:tabLst>
      <a:defRPr sz="1600" u="none" strike="noStrike" kern="1200" cap="none" spc="150" normalizeH="0" baseline="0">
        <a:solidFill><a:schemeClr val="tx1"><a:lumMod val="65000"/><a:lumOff val="35000"/></a:schemeClr></a:solidFill>
        <a:uFillTx/>
        <a:latin typeface="+mn-lt"/>
        <a:ea typeface="+mn-ea"/>
        <a:cs typeface="+mn-cs"/>
      </a:defRPr>
    </a:lvl2pPr>
    <a:lvl3pPr marL="1143000" indent="-228600" algn="l" defTabSz="914400" rtl="0" eaLnBrk="1" fontAlgn="auto" latinLnBrk="0" hangingPunct="1">
      <a:lnSpc><a:spcPct val="120000"/></a:lnSpc>
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="600"/></a:spcAft>
      <a:buFont typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>
      <a:buChar char="●"/>
      <a:defRPr sz="1600" u="none" strike="noStrike" kern="1200" cap="none" spc="150" normalizeH="0" baseline="0">
        <a:solidFill><a:schemeClr val="tx1"><a:lumMod val="65000"/><a:lumOff val="35000"/></a:schemeClr></a:solidFill>
        <a:uFillTx/>
        <a:latin typeface="+mn-lt"/>
        <a:ea typeface="+mn-ea"/>
        <a:cs typeface="+mn-cs"/>
      </a:defRPr>
    </a:lvl3pPr>
    <a:lvl4pPr marL="1600200" indent="-228600" algn="l" defTabSz="914400" rtl="0" eaLnBrk="1" fontAlgn="auto" latinLnBrk="0" hangingPunct="1">
      <a:lnSpc><a:spcPct val="120000"/></a:lnSpc>
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="300"/></a:spcAft>
      <a:buFont typeface="Wingdings" panose="05000000000000000000" charset="0"/>
      <a:buChar char=""/>
      <a:defRPr sz="1400" u="none" strike="noStrike" kern="1200" cap="none" spc="150" normalizeH="0" baseline="0">
        <a:solidFill><a:schemeClr val="tx1"><a:lumMod val="65000"/><a:lumOff val="35000"/></a:schemeClr></a:solidFill>
        <a:uFillTx/>
        <a:latin typeface="+mn-lt"/>
        <a:ea typeface="+mn-ea"/>
        <a:cs typeface="+mn-cs"/>
      </a:defRPr>
    </a:lvl4pPr>
    <a:lvl5pPr marL="2057400" indent="-228600" algn="l" defTabSz="914400" rtl="0" eaLnBrk="1" fontAlgn="auto" latinLnBrk="0" hangingPunct="1">
      <a:lnSpc><a:spcPct val="120000"/></a:lnSpc>
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="300"/></a:spcAft>
      <a:buFont typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>
      <a:buChar char="•"/>
      <a:defRPr sz="1400" u="none" strike="noStrike" kern="1200" cap="none" spc="150" normalizeH="0" baseline="0">
        <a:solidFill><a:schemeClr val="tx1"><a:lumMod val="65000"/><a:lumOff val="35000"/></a:schemeClr></a:solidFill>
        <a:uFillTx/>
        <a:latin typeface="+mn-lt"/>
        <a:ea typeface="+mn-ea"/>
        <a:cs typeface="+mn-cs"/>
      </a:defRPr>
    </a:lvl5pPr>
    <a:lvl6pPr marL="2514600" indent="-228600" algn="l" defTabSz="914400" rtl="0" eaLnBrk="1" latinLnBrk="0" hangingPunct="1">
      <a:lnSpc><a:spcPct val="90000"/></a:lnSpc>
      <a:spcBef><a:spcPts val="500"/></a:spcBef>
      <a:buFont typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>
      <a:buChar char="•"/>
      <a:defRPr sz="1800" kern="1200">
        <a:solidFill><a:schemeClr val="tx1"/></a:solidFill>
        <a:latin typeface="+mn-lt"/>
        <a:ea typeface="+mn-ea"/>
        <a:cs typeface="+mn-cs"/>
      </a:defRPr>
    </a:lvl6pPr>
    <a:lvl7pPr marL="2971800" indent="-228600" algn="l" defTabSz="914400" rtl="0" eaLnBrk="1" latinLnBrk="0" hangingPunct="1">
      <a:lnSpc><a:spcPct val="90000"/></a:lnSpc>
      <a:spcBef><a:spcPts val="500"/></a:spcBef>
      <a:buFont typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>
      <a:buChar char="•"/>
      <a:defRPr sz="1800" kern="1200">
        <a:solidFill><a:schemeClr val="tx1"/></a:solidFill>
        <a:latin typeface="+mn-lt"/>
        <a:ea typeface="+mn-ea"/>
        <a:cs typeface="+mn-cs"/>
      </a:defRPr>
    </a:lvl7pPr>
    <a:lvl8pPr marL="3429000" indent="-228600" algn="l" defTabSz="914400" rtl="0" eaLnBrk="1" latinLnBrk="0" hangingPunct="1">
      <a:lnSpc><a:spcPct val="90000"/></a:lnSpc>
      <a:spcBef><a:spcPts val="500"/></a:spcBef>
      <a:buFont typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>
      <a:buChar char="•"/>
      <a:defRPr sz="1800" kern="1200">
        <a:solidFill><a:schemeClr val="tx1"/></a:solidFill>
        <a:latin typeface="+mn-lt"/>
        <a:ea typeface="+mn-ea"/>
        <a:cs typeface="+mn-cs"/>
      </a:defRPr>
    </a:lvl8pPr>
    <a:lvl9pPr marL="3886200" indent="-228600" algn="l" defTabSz="914400" rtl="0" eaLnBrk="1" latinLnBrk="0" hangingPunct="1">
      <a:lnSpc><a:spcPct val="90000"/></a:lnSpc>
      <a:spcBef><a:spcPts val="500"/></a:spcBef>
      <a:buFont typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>
      <a:buChar char="•"/>
      <a:defRPr sz="1800" kern="1200">
        <a:solidFill><a:schemeClr val="tx1"/></a:solidFill>
        <a:latin typeface="+mn-lt"/>
        <a:ea typeface="+mn-ea"/>
        <a:cs typeface="+mn-cs"/>
      </a:defRPr>
    </a:lvl9pPr>
  </a:lstStyle>"""


def generate_ppt(topics: List[Topic], config: dict, output_path: str = None) -> str:
    """生成完整的填空题PPT"""
    if output_path is None:
        output_path = config["output_path"]

    template_path = config["template_ppt"]
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"模板PPT不存在: {template_path}")

    prs = Presentation(template_path)

    # 提取标题页图片
    title_image_blob = _extract_title_image(prs)

    # 清空已有幻灯片
    _clear_slides(prs)

    # 对所有主题进行分页
    pages = paginate_topics(topics, config)

    # 按主题分组生成
    current_topic_name = None
    for topic_name, page_paras in pages:
        if topic_name != current_topic_name:
            _add_title_slide(prs, topic_name, config, title_image_blob)
            current_topic_name = topic_name
        _add_content_slide(prs, topic_name, page_paras, config)

    # 保存
    prs.save(output_path)
    print(f"[OK] PPT saved to: {output_path}")
    print(f"     Total slides: {len(prs.slides)}")

    return output_path


def _extract_title_image(prs: Presentation):
    """从模板PPT的第一张幻灯片提取图片"""
    if len(prs.slides) == 0:
        return None
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.shape_type == 13:
            try:
                return shape.image.blob
            except Exception:
                pass
    return None


def _clear_slides(prs: Presentation):
    """清空已有幻灯片"""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        sldIdLst.remove(sldId)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass


def _create_shape_from_xml(slide, shape_xml_str: str):
    """从XML字符串创建shape并添加到slide"""
    spTree = slide.shapes._spTree
    new_el = etree.fromstring(shape_xml_str)
    spTree.append(new_el)
    return new_el


def _add_title_slide(prs: Presentation, topic_name: str, config: dict, image_blob=None):
    """添加标题幻灯片（使用占位符结构）"""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # 清空布局自带的占位符
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # 添加背景
    _add_background(slide)

    # 添加标题占位符（使用XML直接构建，匹配参考PPT结构）
    title_xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="2" name="标题 1"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="ctrTitle"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr/>
  <p:txBody>
    <a:bodyPr/>
    <a:p>
      <a:r>
        <a:rPr lang="zh-CN" altLang="zh-CN"/>
        <a:t>{_escape_xml(topic_name)}</a:t>
      </a:r>
      <a:endParaRPr lang="zh-CN" altLang="zh-CN"/>
    </a:p>
  </p:txBody>
</p:sp>"""
    _create_shape_from_xml(slide, title_xml)

    # 添加图片
    if image_blob:
        from io import BytesIO
        slide.shapes.add_picture(
            BytesIO(image_blob),
            Emu(IMG_LEFT), Emu(IMG_TOP),
            Emu(IMG_W), Emu(IMG_H),
        )


def _add_content_slide(prs: Presentation, topic_name: str, paragraphs: List[Paragraph], config: dict):
    """添加内容幻灯片（使用占位符结构，匹配参考PPT）"""
    layout_idx = min(1, len(prs.slide_layouts) - 1)
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    # 清空布局自带的占位符
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # 添加背景
    _add_background(slide)

    # 1. 添加标题占位符
    title_xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="2" name="标题 1"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="title"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr/>
  <p:txBody>
    <a:bodyPr/>
    <a:p>
      <a:r>
        <a:rPr lang="zh-CN" altLang="en-US"/>
        <a:t>{_escape_xml(topic_name)}</a:t>
      </a:r>
      <a:endParaRPr lang="zh-CN" altLang="en-US"/>
    </a:p>
  </p:txBody>
</p:sp>"""
    _create_shape_from_xml(slide, title_xml)

    # 2. 添加基础层文本框（使用占位符结构）
    base_shape_id = 3
    base_xml = _build_base_layer_xml(topic_name, paragraphs, base_shape_id, config)
    _create_shape_from_xml(slide, base_xml)

    # 3. 收集本页所有挖空词
    all_blanks = []
    for para in paragraphs:
        all_blanks.extend(para.blanks)

    if not all_blanks:
        return

    # 去重
    seen = set()
    unique_blanks = []
    for b in all_blanks:
        if b not in seen:
            seen.add(b)
            unique_blanks.append(b)

    # 4. 为每个挖空词创建答案层 shape（克隆基础层结构）
    answer_shape_specs = []
    next_shape_id = 4

    for blank in unique_blanks:
        # 找到挖空词所在的段落索引（用 blanks 精确匹配，兼容 manual 和 AI 模式）
        para_idx = -1
        for i, para in enumerate(paragraphs):
            if blank in para.blanks:
                para_idx = i
                break
        
        if para_idx == -1:
            # 回退：用 clean_text 搜索
            for i, para in enumerate(paragraphs):
                if blank in para.clean_text():
                    para_idx = i
                    break

        if para_idx == -1:
            print(f"[WARN] 挖空词 '{blank}' 不在本页文本中，跳过")
            continue

        answer_xml = _build_answer_layer_xml(topic_name, paragraphs, blank, next_shape_id, config, target_para_idx=para_idx)
        _create_shape_from_xml(slide, answer_xml)
        answer_shape_specs.append((next_shape_id, para_idx))
        next_shape_id += 1

    # 5. 添加动画
    if answer_shape_specs:
        add_click_fade_animation(slide, answer_shape_specs, config)


def _build_base_layer_xml(topic_name: str, paragraphs: List[Paragraph], shape_id: int, config: dict) -> str:
    """构建基础层占位符的完整XML"""
    paragraphs_xml = ""
    for para in paragraphs:
        # 使用原始文本（含{{}}标记）精确定位挖空词
        segments = split_text_by_blanks(para.text, para.blanks)

        runs_xml = ""
        for segment_text, is_blank in segments:
            if not segment_text:
                continue
            # 处理 \n：按换行拆分，用 <a:br/> 连接
            parts = segment_text.split('\n')
            for pi, part in enumerate(parts):
                if pi > 0:
                    runs_xml += '<a:br/>'
                if not part:
                    continue
                escaped = _escape_xml(part)
                if is_blank:
                    # 挖空词：透明文字 + 黑色下划线
                    runs_xml += f"""<a:r><a:rPr lang="zh-CN" altLang="en-US" sz="{FONT_SIZE}" u="sng"><a:noFill/><a:uFill><a:solidFill><a:schemeClr val="tx1"/></a:solidFill></a:uFill></a:rPr><a:t>{escaped}</a:t></a:r>"""
                else:
                    # 非挖空词：黑色文字（继承layout，不显式设置颜色）
                    runs_xml += f"""<a:r><a:rPr lang="zh-CN" altLang="en-US" sz="{FONT_SIZE}"/><a:t>{escaped}</a:t></a:r>"""

        paragraphs_xml += f"""<a:p><a:pPr marL="0" indent="0"><a:buNone/></a:pPr>{runs_xml}<a:endParaRPr lang="zh-CN" altLang="en-US" sz="{FONT_SIZE}"/></a:p>"""

    xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="内容占位符 2"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph idx="1"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr/>
  <p:txBody>
    <a:bodyPr><a:noAutofit/></a:bodyPr>
    {paragraphs_xml}
  </p:txBody>
</p:sp>"""
    return xml


def _build_answer_layer_xml(topic_name: str, paragraphs: List[Paragraph], blank: str, shape_id: int, config: dict, target_para_idx: int = -1) -> str:
    """构建答案层shape的完整XML（克隆基础层结构，仅高亮当前blank）

    Args:
        target_para_idx: 答案文字所在的段落索引。只有该段落中的 {{blank}} 才会高亮为蓝色。
                         其他段落中的同名文本保持透明，避免动画未隐藏的蓝字泄露。
    """
    paragraphs_xml = ""
    for para_idx, para in enumerate(paragraphs):
        # 使用原始文本（含{{}}标记）精确定位挖空词
        segments = split_text_by_blanks(para.text, para.blanks)

        runs_xml = ""
        for segment_text, is_blank in segments:
            if not segment_text:
                continue
            # 处理 \n：按换行拆分，用 <a:br/> 连接
            parts = segment_text.split('\n')
            for pi, part in enumerate(parts):
                if pi > 0:
                    runs_xml += '<a:br/>'
                if not part:
                    continue
                escaped = _escape_xml(part)
                if is_blank and segment_text == blank and (target_para_idx < 0 or para_idx == target_para_idx):
                    # 答案词：蓝色（仅在目标段落高亮）
                    runs_xml += f"""<a:r><a:rPr lang="zh-CN" altLang="en-US" sz="{FONT_SIZE}"><a:solidFill><a:schemeClr val="accent1"><a:lumMod val="60000"/><a:lumOff val="40000"/></a:schemeClr></a:solidFill></a:rPr><a:t>{escaped}</a:t></a:r>"""
                else:
                    # 非答案词或非目标段落：透明
                    runs_xml += f"""<a:r><a:rPr lang="zh-CN" altLang="en-US" sz="{FONT_SIZE}"><a:noFill/></a:rPr><a:t>{escaped}</a:t></a:r>"""

        paragraphs_xml += f"""<a:p><a:pPr marL="0" indent="0"><a:buNone/></a:pPr>{runs_xml}<a:endParaRPr lang="zh-CN" altLang="en-US" sz="{FONT_SIZE}"><a:noFill/></a:endParaRPr></a:p>"""

    xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="内容占位符 2"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{CONTENT_LEFT}" y="{CONTENT_TOP}"/>
      <a:ext cx="{CONTENT_W}" cy="{CONTENT_H}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr vert="horz" lIns="{MARGIN_L}" tIns="{MARGIN_T}" rIns="{MARGIN_R}" bIns="{MARGIN_B}" rtlCol="0"><a:noAutofit/></a:bodyPr>
    {ANSWER_LST_STYLE}
    {paragraphs_xml}
  </p:txBody>
</p:sp>"""
    return xml


def _add_background(slide):
    """添加白色背景"""
    bg_xml = """<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><p:bgRef idx="1001"><a:schemeClr val="bg1"/></p:bgRef></p:bg>"""
    bg_el = etree.fromstring(bg_xml)
    # 插入到 spTree 之前（bg 必须在 spTree 之前）
    spTree = slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    slide._element.find(qn('p:cSld')).insert(0, bg_el)


def _escape_xml(text: str) -> str:
    """转义XML特殊字符并过滤非法控制字符"""
    if not text:
        return ''
    # 过滤 XML 1.0 非法控制字符 (0x00-0x08, 0x0B, 0x0C, 0x0E-0x1F)
    # 保留: 0x09(\t), 0x0A(\n), 0x0D(\r)
    import re as _re
    text = _re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
    # 转义 XML 特殊字符
    text = text.replace('&', '&amp;')
    text = text.replace('<', '&lt;')
    text = text.replace('>', '&gt;')
    text = text.replace('"', '&quot;')
    text = text.replace("'", '&apos;')
    return text


def save_memory_draft(output_path: str, memory_dir: str):
    """保存PPT草稿到Memory"""
    import shutil
    os.makedirs(memory_dir, exist_ok=True)
    dst = os.path.join(memory_dir, '04_draft.pptx')
    shutil.copy2(output_path, dst)
    return dst

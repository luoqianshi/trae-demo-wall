"""自动验证器模块 v2 - 适配新的占位符结构

解析生成的PPT，运行验证用例，输出0/1信号
"""
import os
import json
from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from lxml import etree
from pptx.oxml.ns import qn


def validate_pptx(pptx_path: str, config: dict) -> Dict:
    """验证生成的PPT"""
    results = []

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        results.append({"tc_id": "TC-01", "passed": False, "reason": f"无法打开PPT: {e}", "fix_step": 4})
        return _build_result(results)

    # TC-01: 基本结构验证
    passed = len(prs.slides) > 0
    results.append({"tc_id": "TC-01", "passed": passed, "reason": "" if passed else "幻灯片数量为0", "fix_step": 4})

    # TC-02: 幻灯片尺寸
    expected_w = Inches(config["slide_width_inch"])
    expected_h = Inches(config["slide_height_inch"])
    passed = (abs(prs.slide_width - expected_w) < Inches(0.1) and abs(prs.slide_height - expected_h) < Inches(0.1))
    results.append({"tc_id": "TC-02", "passed": passed, "reason": "" if passed else f"尺寸不符", "fix_step": 4})

    # TC-03: 检查标题页和内容页
    has_title = len(prs.slides) > 0
    has_content = len(prs.slides) > 1
    passed = has_title and has_content
    results.append({"tc_id": "TC-03", "passed": passed, "reason": "" if passed else "缺少标题页或内容页", "fix_step": 4})

    # TC-04: 检查挖空词格式（基础层 noFill + underline）
    base_layer_ok = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is not None:
                        has_nofill = rPr.find(qn('a:noFill')) is not None
                        has_ufill = rPr.find(qn('a:uFill')) is not None
                        has_underline = rPr.get('u') == 'sng'
                        if has_nofill and has_ufill and has_underline:
                            base_layer_ok = True
                            break
                if base_layer_ok:
                    break
            if base_layer_ok:
                break
        if base_layer_ok:
            break
    passed = base_layer_ok
    results.append({"tc_id": "TC-04", "passed": passed, "reason": "" if passed else "未找到基础层挖空格式", "fix_step": 4})

    # TC-05: 检查答案层格式（accent1 + lumMod + lumOff）
    answer_layer_ok = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is not None:
                        solid_fill = rPr.find(qn('a:solidFill'))
                        if solid_fill is not None:
                            scheme_clr = solid_fill.find(qn('a:schemeClr'))
                            if scheme_clr is not None and scheme_clr.get('val') == 'accent1':
                                lum_mod = scheme_clr.find(qn('a:lumMod'))
                                lum_off = scheme_clr.find(qn('a:lumOff'))
                                if lum_mod is not None and lum_off is not None:
                                    answer_layer_ok = True
                                    break
                if answer_layer_ok:
                    break
            if answer_layer_ok:
                break
        if answer_layer_ok:
            break
    passed = answer_layer_ok
    results.append({"tc_id": "TC-05", "passed": passed, "reason": "" if passed else "未找到答案层蓝色格式", "fix_step": 4})

    # TC-06: 检查动画
    animation_ok = False
    for slide in prs.slides:
        timing = slide._element.find(qn('p:timing'))
        if timing is not None:
            set_elements = timing.findall('.//' + qn('p:set'))
            anim_effects = timing.findall('.//' + qn('p:animEffect'))
            if set_elements and anim_effects:
                for s in set_elements:
                    attr_names = s.findall('.//' + qn('p:attrName'))
                    for attr in attr_names:
                        if attr.text == 'style.visibility':
                            animation_ok = True
                            break
                    if animation_ok:
                        break
            if animation_ok:
                break
    passed = animation_ok
    results.append({"tc_id": "TC-06", "passed": passed, "reason": "" if passed else "未找到点击淡入动画", "fix_step": 4})

    # TC-07: 检查 bldP 是否有 build="p"
    bldp_ok = False
    for slide in prs.slides:
        timing = slide._element.find(qn('p:timing'))
        if timing is not None:
            bldps = timing.findall('.//' + qn('p:bldP'))
            for b in bldps:
                if b.get('build') == 'p':
                    bldp_ok = True
                    break
        if bldp_ok:
            break
    passed = bldp_ok
    results.append({"tc_id": "TC-07", "passed": passed, "reason": "" if passed else "bldP缺少build='p'", "fix_step": 4})

    # TC-08: 检查答案层shape有spLocks
    splocks_ok = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # 检查是否有多个shape（基础层+答案层）
                sp = shape._element
                cNvSpPr = sp.find('.//' + qn('p:cNvSpPr'))
                if cNvSpPr is not None:
                    spLocks = cNvSpPr.find(qn('a:spLocks'))
                    if spLocks is not None:
                        splocks_ok = True
                        break
        if splocks_ok:
            break
    passed = splocks_ok
    results.append({"tc_id": "TC-08", "passed": passed, "reason": "" if passed else "shape缺少spLocks", "fix_step": 4})

    # TC-09: 检查 pRg 指向包含蓝色答案文字的段落
    prg_ok = True
    prg_reason = ""
    for slide in prs.slides:
        timing = slide._element.find(qn('p:timing'))
        if timing is None:
            continue
        
        # 收集每个 spid 对应的 pRg 值
        spTgts = timing.findall('.//' + qn('p:spTgt'))
        spid_to_prg = {}
        for spTgt in spTgts:
            spid = spTgt.get('spid')
            txEl = spTgt.find(qn('p:txEl'))
            if txEl is not None:
                pRg = txEl.find(qn('p:pRg'))
                if pRg is not None:
                    st = int(pRg.get('st', '0'))
                    # 只记录第一次出现的 pRg（set 和 animEffect 都有，值相同）
                    if spid not in spid_to_prg:
                        spid_to_prg[spid] = st
        
        # 对每个答案 shape，检查 pRg 是否指向蓝色文字所在段落
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            sp = shape._element
            cNvPr = sp.find('.//' + qn('p:cNvPr'))
            if cNvPr is None:
                continue
            sid = cNvPr.get('id')
            if sid not in spid_to_prg:
                continue
            
            # 找到蓝色文字所在段落索引
            blue_para_idx = -1
            for i, para in enumerate(shape.text_frame.paragraphs):
                for run in para.runs:
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is not None:
                        sf = rPr.find(qn('a:solidFill'))
                        if sf is not None:
                            sc = sf.find(qn('a:schemeClr'))
                            if sc is not None and sc.get('val') == 'accent1':
                                blue_para_idx = i
                                break
                if blue_para_idx >= 0:
                    break
            
            if blue_para_idx >= 0:
                expected_prg = blue_para_idx
                actual_prg = spid_to_prg[sid]
                if actual_prg != expected_prg:
                    prg_ok = False
                    prg_reason = f"spid={sid}: pRg st={actual_prg} 但蓝色文字在段落{expected_prg}"
                    break
            else:
                # 答案 shape 不含任何蓝色文字 — 严重错误
                prg_ok = False
                prg_reason = f"spid={sid}: 答案shape不含任何蓝色文字"
                break
        
        if not prg_ok:
            break
    
    results.append({"tc_id": "TC-09", "passed": prg_ok, "reason": prg_reason, "fix_step": 4})

    # TC-10: 检查每个答案 shape 都有对应的动画（bldP + set + animEffect）
    anim_ok = True
    anim_reason = ""
    for slide in prs.slides:
        timing = slide._element.find(qn('p:timing'))
        if timing is None:
            continue
        
        # 收集答案 shape ID（无 ph 且有蓝色文字的 shape）
        answer_spids = set()
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            sp = shape._element
            ph = sp.find('.//' + qn('p:ph'))
            if ph is not None:
                continue  # 跳过 placeholder
            # 检查是否有蓝色文字
            has_blue = False
            for run in shape.text_frame.paragraphs:
                for r in run.runs:
                    rPr = r._r.find(qn('a:rPr'))
                    if rPr is not None:
                        sf = rPr.find(qn('a:solidFill'))
                        if sf is not None:
                            sc = sf.find(qn('a:schemeClr'))
                            if sc is not None and sc.get('val') == 'accent1':
                                has_blue = True
                                break
                    if has_blue:
                        break
                if has_blue:
                    break
            if has_blue:
                cNvPr = sp.find('.//' + qn('p:cNvPr'))
                if cNvPr is not None:
                    answer_spids.add(cNvPr.get('id'))
        
        if not answer_spids:
            continue
        
        # 收集有动画的 shape ID
        bldP_spids = set()
        set_spids = set()
        anim_spids = set()
        
        for bldP in timing.findall('.//' + qn('p:bldP')):
            bldP_spids.add(bldP.get('spid'))
        
        for s in timing.findall('.//' + qn('p:set')):
            tgtEl = s.find('.//' + qn('p:tgtEl'))
            if tgtEl is not None:
                spTgt = tgtEl.find(qn('p:spTgt'))
                if spTgt is not None:
                    set_spids.add(spTgt.get('spid'))
        
        for anim in timing.findall('.//' + qn('p:animEffect')):
            tgtEl = anim.find('.//' + qn('p:tgtEl'))
            if tgtEl is not None:
                spTgt = tgtEl.find(qn('p:spTgt'))
                if spTgt is not None:
                    anim_spids.add(spTgt.get('spid'))
        
        # 检查每个答案 shape 是否都有完整的动画
        for spid in answer_spids:
            if spid not in bldP_spids:
                anim_ok = False
                anim_reason = f"spid={spid}: 缺少bldP动画"
                break
            if spid not in set_spids:
                anim_ok = False
                anim_reason = f"spid={spid}: 缺少set(visibility)动画"
                break
            if spid not in anim_spids:
                anim_ok = False
                anim_reason = f"spid={spid}: 缺少animEffect(fade)动画"
                break
        
        if not anim_ok:
            break
    
    results.append({"tc_id": "TC-10", "passed": anim_ok, "reason": anim_reason, "fix_step": 4})

    return _build_result(results)


def _build_result(results: List[Dict]) -> Dict:
    all_passed = all(r["passed"] for r in results)
    fix_step = None
    if not all_passed:
        failed = [r for r in results if not r["passed"]]
        fix_step = min(r["fix_step"] for r in failed) if failed else None
    return {"all_passed": all_passed, "results": results, "fix_target_step": fix_step}


def save_memory(validation: Dict, memory_dir: str):
    os.makedirs(memory_dir, exist_ok=True)
    path = os.path.join(memory_dir, '05_validation.json')
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(validation, f, ensure_ascii=False, indent=2)
    return path

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

TEMPLATES = {
    "education": {
        "name": "教育汇报模板",
        "primary_color": RGBColor(102, 126, 234),
        "secondary_color": RGBColor(118, 75, 162),
        "accent_color": RGBColor(40, 167, 69),
        "title_font": "微软雅黑",
        "body_font": "微软雅黑",
        "bg_color": RGBColor(248, 249, 255)
    },
    "student": {
        "name": "学生展示模板",
        "primary_color": RGBColor(255, 107, 107),
        "secondary_color": RGBColor(238, 90, 36),
        "accent_color": RGBColor(255, 193, 7),
        "title_font": "微软雅黑",
        "body_font": "微软雅黑",
        "bg_color": RGBColor(255, 248, 240)
    },
    "business": {
        "name": "数据报告模板",
        "primary_color": RGBColor(33, 37, 41),
        "secondary_color": RGBColor(60, 66, 82),
        "accent_color": RGBColor(52, 152, 219),
        "title_font": "微软雅黑",
        "body_font": "微软雅黑",
        "bg_color": RGBColor(255, 255, 255)
    }
}


class PPTGenerator:
    def __init__(self, output_dir):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def create_presentation(self, template_name="education"):
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        return prs

    def get_template(self, template_name):
        return TEMPLATES.get(template_name, TEMPLATES["education"])

    def add_title_slide(self, prs, title, subtitle="", template_name="education"):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        template = self.get_template(template_name)
        
        self._apply_background(slide, template)
        
        title_box = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(2.5), 
            width=Inches(14), height=Inches(2)
        )
        title_text = title_box.text_frame.add_paragraph()
        title_text.text = title
        title_text.font.size = Pt(52)
        title_text.font.bold = True
        title_text.font.color.rgb = template["primary_color"]
        title_text.font.name = template["title_font"]
        title_text.alignment = PP_ALIGN.CENTER
        title_text.space_after = Pt(20)
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                left=Inches(2), top=Inches(4.5), 
                width=Inches(12), height=Inches(1.5)
            )
            subtitle_text = subtitle_box.text_frame.add_paragraph()
            subtitle_text.text = subtitle
            subtitle_text.font.size = Pt(24)
            subtitle_text.font.color.rgb = template["secondary_color"]
            subtitle_text.font.name = template["body_font"]
            subtitle_text.alignment = PP_ALIGN.CENTER
        
        self._add_title_banner(slide, template)
        
        return slide

    def add_content_slide(self, prs, title, content="", template_name="education"):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        template = self.get_template(template_name)
        
        self._apply_background(slide, template)
        self._add_section_banner(slide, template)
        
        title_box = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(0.5), 
            width=Inches(14), height=Inches(1)
        )
        title_text = title_box.text_frame.add_paragraph()
        title_text.text = title
        title_text.font.size = Pt(36)
        title_text.font.bold = True
        title_text.font.color.rgb = template["primary_color"]
        title_text.font.name = template["title_font"]
        
        if content:
            content_box = slide.shapes.add_textbox(
                left=Inches(1), top=Inches(2), 
                width=Inches(14), height=Inches(6)
            )
            content_text = content_box.text_frame.add_paragraph()
            content_text.text = content
            content_text.font.size = Pt(20)
            content_text.font.color.rgb = template["secondary_color"]
            content_text.font.name = template["body_font"]
            content_text.line_spacing = Pt(28)
        
        return slide

    def add_list_slide(self, prs, title, items=None, template_name="education"):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        template = self.get_template(template_name)
        
        self._apply_background(slide, template)
        self._add_section_banner(slide, template)
        
        title_box = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(0.5), 
            width=Inches(14), height=Inches(1)
        )
        title_text = title_box.text_frame.add_paragraph()
        title_text.text = title
        title_text.font.size = Pt(36)
        title_text.font.bold = True
        title_text.font.color.rgb = template["primary_color"]
        title_text.font.name = template["title_font"]
        
        if items and isinstance(items, list):
            list_box = slide.shapes.add_textbox(
                left=Inches(1.5), top=Inches(2), 
                width=Inches(13), height=Inches(6)
            )
            text_frame = list_box.text_frame
            
            for i, item in enumerate(items):
                paragraph = text_frame.add_paragraph()
                paragraph.text = item
                paragraph.font.size = Pt(22)
                paragraph.font.color.rgb = template["secondary_color"]
                paragraph.font.name = template["body_font"]
                paragraph.font.bold = True if i == 0 else False
                paragraph.space_after = Pt(12)
        
        return slide

    def add_image_slide(self, prs, title, image_path=None, description="", template_name="education"):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        template = self.get_template(template_name)
        
        self._apply_background(slide, template)
        
        title_box = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(0.5), 
            width=Inches(14), height=Inches(1)
        )
        title_text = title_box.text_frame.add_paragraph()
        title_text.text = title
        title_text.font.size = Pt(36)
        title_text.font.bold = True
        title_text.font.color.rgb = template["primary_color"]
        title_text.font.name = template["title_font"]
        
        if image_path and os.path.exists(image_path):
            img_left = Inches(1.5)
            img_top = Inches(1.8)
            img_width = Inches(13)
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
        
        if description:
            desc_box = slide.shapes.add_textbox(
                left=Inches(1), top=Inches(7.5), 
                width=Inches(14), height=Inches(1)
            )
            desc_text = desc_box.text_frame.add_paragraph()
            desc_text.text = description
            desc_text.font.size = Pt(18)
            desc_text.font.color.rgb = template["secondary_color"]
            desc_text.font.name = template["body_font"]
            desc_text.alignment = PP_ALIGN.CENTER
        
        return slide

    def add_chart_slide(self, prs, title, chart_type="bar", data=None, categories=None, template_name="education"):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        template = self.get_template(template_name)
        
        self._apply_background(slide, template)
        self._add_section_banner(slide, template)
        
        title_box = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(0.5), 
            width=Inches(14), height=Inches(1)
        )
        title_text = title_box.text_frame.add_paragraph()
        title_text.text = title
        title_text.font.size = Pt(36)
        title_text.font.bold = True
        title_text.font.color.rgb = template["primary_color"]
        title_text.font.name = template["title_font"]
        
        if data and categories:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            
            chart_data = CategoryChartData()
            chart_data.categories = categories
            
            if isinstance(data, dict):
                for series_name, values in data.items():
                    chart_data.add_series(series_name, values)
            
            x, y, cx, cy = Inches(1), Inches(2), Inches(14), Inches(6)
            
            if chart_type == "pie":
                chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
            elif chart_type == "line":
                chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
            else:
                chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
            
            chart.has_legend = True
        
        return slide

    def save_presentation(self, prs, filename=None):
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{timestamp}.pptx"
        
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        
        return filepath

    def _apply_background(self, slide, template):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = template["bg_color"]

    def _add_title_banner(self, slide, template):
        banner_height = Emu(int(Pt(12).emu))
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(0),
            top=Inches(1.8),
            width=Inches(16),
            height=banner_height
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = template["accent_color"]
        banner.line.fill.background()

    def _add_section_banner(self, slide, template):
        banner_height = Emu(int(Pt(8).emu))
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(0),
            top=Inches(1.3),
            width=Inches(16),
            height=banner_height
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = template["accent_color"]
        banner.line.fill.background()

    def get_available_templates(self):
        return [
            {"id": key, "name": template["name"]}
            for key, template in TEMPLATES.items()
        ]

    def generate_from_slides(self, slides_data, template_name="education"):
        prs = self.create_presentation(template_name)
        
        for slide_data in slides_data:
            if hasattr(slide_data, 'dict'):
                slide_dict = slide_data.dict()
            else:
                slide_dict = slide_data
            
            slide_type = slide_dict.get("type", "content")
            title = slide_dict.get("title", "")
            
            if slide_type == "title":
                subtitle = slide_dict.get("subtitle", "")
                self.add_title_slide(prs, title, subtitle, template_name)
            elif slide_type == "list":
                items = slide_dict.get("items", [])
                self.add_list_slide(prs, title, items, template_name)
            elif slide_type == "image":
                image_path = slide_dict.get("image_path", "")
                description = slide_dict.get("description", "")
                self.add_image_slide(prs, title, image_path, description, template_name)
            elif slide_type == "chart":
                chart_type = slide_dict.get("chart_type", "bar")
                data = slide_dict.get("data", {})
                categories = slide_dict.get("categories", [])
                self.add_chart_slide(prs, title, chart_type, data, categories, template_name)
            else:
                content = slide_dict.get("content", "")
                self.add_content_slide(prs, title, content, template_name)
        
        return prs